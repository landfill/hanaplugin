#!/usr/bin/env python3
"""
태깅 정의서 마크다운 데이터를 읽어 .pptx 생성.
사용법: python3 generate_pptx.py <데이터.md> <출력.pptx> [템플릿.pptx]

템플릿 사용: 스킬 폴더 template/ 태깅정의서_템플릿.pptx 의 표 서식을 유지합니다.
- 태깅 상세 슬라이드(인덱스 3)를 XML 복제해 이벤트 수만큼 생성합니다.
- 복제된 슬라이드에는 제목·표·블릿·"태깅 항목" 등 모든 서식이 그대로 유지되며, 데이터만 교체합니다.
"""
import copy
import sys
from pathlib import Path

from lxml import etree  # python-pptx 필수 의존성이므로 항상 사용 가능

SCRIPT_DIR = Path(__file__).resolve().parent
DEFAULT_TEMPLATE = SCRIPT_DIR.parent / "template" / "태깅정의서_템플릿.pptx"


def _find_project_root(start_path: Path) -> Path:
    """start_path에서 위로 올라가며 프로젝트 루트를 탐지. .git, .cursor-plugin, .cursor 중 하나가 있는 디렉토리를 반환."""
    current = start_path.resolve()
    if current.is_file():
        current = current.parent
    for _ in range(20):
        if (current / ".git").exists() or (current / ".cursor-plugin").exists() or (current / ".cursor").exists():
            return current
        parent = current.parent
        if parent == current:
            break
        current = parent
    return start_path.resolve().parent


def _resolve_image_path(img_path_raw: str, md_dir: Path, project_root: Path) -> "Path | None":
    """이미지 경로를 여러 기준으로 탐색해 실제 존재하는 절대경로를 반환."""
    if not img_path_raw:
        return None
    candidates = [
        md_dir / img_path_raw,           # 1) 마크다운 파일 기준
        project_root / img_path_raw,     # 2) 프로젝트 루트 기준
        Path(img_path_raw),              # 3) 절대경로 또는 CWD 기준
    ]
    for c in candidates:
        resolved = c.resolve()
        if resolved.exists() and resolved.is_file():
            return resolved
    return None


def _get_paragraph_font(paragraph):
    """단락의 폰트 정보 반환 (size, bold, name). runs[0] 우선, 없으면 paragraph.font."""
    try:
        if paragraph.runs:
            r = paragraph.runs[0]
            return {
                "size": getattr(r.font, "size", None),
                "bold": getattr(r.font, "bold", None),
                "name": getattr(r.font, "name", None),
            }
        return {
            "size": getattr(paragraph.font, "size", None),
            "bold": getattr(paragraph.font, "bold", None),
            "name": getattr(paragraph.font, "name", None),
        }
    except Exception:
        return {"size": None, "bold": None, "name": None}


def _apply_paragraph_font(paragraph, font_dict):
    """단락의 첫 run(또는 단락)에 폰트 적용. 텍스트 설정 후 호출."""
    if not font_dict:
        return
    try:
        if paragraph.runs:
            r = paragraph.runs[0]
            if font_dict.get("size") is not None:
                r.font.size = font_dict["size"]
            if font_dict.get("bold") is not None:
                r.font.bold = font_dict["bold"]
            if font_dict.get("name") is not None:
                r.font.name = font_dict["name"]
        else:
            if font_dict.get("size") is not None:
                paragraph.font.size = font_dict["size"]
            if font_dict.get("bold") is not None:
                paragraph.font.bold = font_dict["bold"]
            if font_dict.get("name") is not None:
                paragraph.font.name = font_dict["name"]
    except Exception:
        pass


def _copy_cell_format(source_cell, target_cell):
    """target_cell에 이미 텍스트가 설정된 상태에서, source_cell의 폰트·정렬을 복사."""
    try:
        if not source_cell.text_frame.paragraphs or not target_cell.text_frame.paragraphs:
            return
        src_para = source_cell.text_frame.paragraphs[0]
        tgt_para = target_cell.text_frame.paragraphs[0]
        font_dict = _get_paragraph_font(src_para)
        _apply_paragraph_font(tgt_para, font_dict)
        if hasattr(src_para, "alignment") and src_para.alignment is not None:
            try:
                tgt_para.alignment = src_para.alignment
            except Exception:
                pass
    except Exception:
        pass


# 태깅 상세: 이미지 유무와 관계없이 태깅항목 표 동일 위치(오른쪽 고정)
TAGGING_TABLE_LEFT_INCH = 4.5
TAGGING_TABLE_TOP_INCH = 1.5
TAGGING_TABLE_WIDTH_INCH = 5.0
WIREFRAME_LEFT_INCH = 0.5
WIREFRAME_TOP_INCH = 1.5
WIREFRAME_WIDTH_INCH = 4.0
WIREFRAME_HEIGHT_INCH = 3.0


def _duplicate_slide(prs, src_slide_idx):
    """템플릿 슬라이드를 XML 통째로 복제해 프레젠테이션 끝에 추가. 이미지는 복사하지 않음(이벤트 루프에서 처리)."""
    src_slide = prs.slides[src_slide_idx]
    slide_layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    new_spTree = new_slide.shapes._spTree
    src_spTree = src_slide.shapes._spTree
    # 레이아웃에서 생긴 기본 shape 제거 (nvGrpSpPr, grpSpPr는 유지)
    for child in list(new_spTree)[2:]:
        new_spTree.remove(child)
    # 원본 shape를 깊은 복사 (이미지(pic)는 제외 — 이벤트별로 별도 처리)
    for child in list(src_spTree)[2:]:
        tag_local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag_local == "pic":
            continue
        new_spTree.append(copy.deepcopy(child))
    return new_slide


def _add_wireframe_placeholder(slide):
    """캡처 없을 때 왼쪽에 와이어프레임 구역(테두리 사각형) 추가."""
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        left = Inches(WIREFRAME_LEFT_INCH)
        top = Inches(WIREFRAME_TOP_INCH)
        width = Inches(WIREFRAME_WIDTH_INCH)
        height = Inches(WIREFRAME_HEIGHT_INCH)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
        )
        shape.fill.background()
        shape.line.color.rgb = RGBColor(0xAD, 0xAD, 0xAD)
        shape.line.width = Pt(0.5)
    except Exception as e:
        print(f"경고: 와이어프레임 추가 실패: {e}", file=sys.stderr)


def _apply_tagging_table_overrides_only(table_shape):
    """태깅 표: 행 높이 균등, 폰트 크기 동일, 헤더 검정+연한 회색 배경, 헤더 외 셀 음영 제거, 셀 세로 가운데."""
    try:
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        tbl = table_shape.table
        nrows = len(tbl.rows)
        ncols = len(tbl.rows[0].cells) if tbl.rows else 0
        if nrows == 0 or ncols == 0:
            return
        # 헤더 행 배경: 연한 회색 (템플릿과 동일)
        HEADER_GRAY = RGBColor(0xE7, 0xE6, 0xE6)
        for r in range(nrows):
            for c in range(ncols):
                try:
                    cell = tbl.cell(r, c)
                    if r == 0:
                        try:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = HEADER_GRAY
                        except Exception:
                            pass
                    else:
                        try:
                            cell.fill.background()
                        except Exception:
                            pass
                except Exception:
                    pass
        # 행 높이 균등
        try:
            total_height = table_shape.height
            row_height_emu = int(total_height / nrows)
            for r in range(nrows):
                tbl.rows[r].height = row_height_emu
        except Exception:
            pass
        # 폰트 크기 동일: 헤더 또는 첫 데이터 셀에서 참조 후 전체 셀에 적용
        ref_size = None
        try:
            if nrows > 1 and tbl.cell(1, 0).text_frame.paragraphs:
                ref_size = _get_paragraph_font(tbl.cell(1, 0).text_frame.paragraphs[0]).get("size")
            if ref_size is None and tbl.cell(0, 0).text_frame.paragraphs:
                ref_size = _get_paragraph_font(tbl.cell(0, 0).text_frame.paragraphs[0]).get("size")
            if ref_size is None:
                ref_size = Pt(9)
        except Exception:
            ref_size = Pt(9)
        for r in range(nrows):
            for c in range(ncols):
                try:
                    cell = tbl.cell(r, c)
                    tf = cell.text_frame
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    try:
                        tf.margin_top = Emu(0)
                        tf.margin_bottom = Emu(0)
                    except Exception:
                        pass
                    for para in tf.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        if ref_size is not None:
                            if para.runs:
                                para.runs[0].font.size = ref_size
                            else:
                                para.font.size = ref_size
                        if r == 0:
                            if para.runs:
                                para.runs[0].font.color.rgb = RGBColor(0, 0, 0)
                            else:
                                para.font.color.rgb = RGBColor(0, 0, 0)
                except Exception:
                    pass
    except Exception as e:
        print(f"경고: 태깅 표 오버라이드 적용 중 오류: {e}", file=sys.stderr)


def _set_text_frame_lines(shape, lines: list, clear_first: bool = True, preserve_font_from_shape=None):
    """텍스트 프레임에 여러 줄 반영. clear_first=True면 기존 단락을 먼저 비운 뒤 채운다.
    preserve_font_from_shape가 있으면 해당 shape의 첫 단락 폰트를 읽어 모든 줄에 적용(개요 본문 등)."""
    if not shape.has_text_frame or not lines:
        return
    tf = shape.text_frame
    font_ref = None
    if preserve_font_from_shape and preserve_font_from_shape.has_text_frame and preserve_font_from_shape.text_frame.paragraphs:
        font_ref = _get_paragraph_font(preserve_font_from_shape.text_frame.paragraphs[0])
    if clear_first:
        for p in tf.paragraphs:
            p.text = ""
    for i, line in enumerate(lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            if i < len(tf.paragraphs):
                tf.paragraphs[i].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0
        if font_ref:
            _apply_paragraph_font(tf.paragraphs[i], font_ref)
    for j in range(len(lines), len(tf.paragraphs)):
        tf.paragraphs[j].text = ""
        if font_ref:
            _apply_paragraph_font(tf.paragraphs[j], font_ref)


def _fill_table_from_template(tbl, headers: list, rows: list, copy_style_from=None):
    """기존 테이블에 헤더+행 채우기. 템플릿 문구는 우리 데이터로 대체하고, 남는 행은 비움.
    copy_style_from(동일 테이블)의 해당 셀 폰트·정렬을 셀마다 복사해 서식 유지.
    모든 셀에 세로 가운데 정렬을 강제 적용."""
    try:
        from pptx.enum.text import MSO_ANCHOR
        nrows_total = len(tbl.rows)
        ncols = len(tbl.rows[0].cells) if tbl.rows else 0
        data_start_col = 0
        if ncols >= 5 and len(headers) == 4:
            data_start_col = 1
        # 먼저 복사할 셀의 서식 수집 (같은 테이블이므로 덮어쓰기 전에 읽음)
        def cell_format_ref(r, c):
            if copy_style_from is None or r >= len(copy_style_from.rows) or c >= len(copy_style_from.rows[0].cells):
                return None, None
            try:
                cell = copy_style_from.cell(r, c)
                if cell.text_frame.paragraphs:
                    para = cell.text_frame.paragraphs[0]
                    font_ref = _get_paragraph_font(para)
                    align = getattr(para, "alignment", None)
                    return font_ref, align
            except Exception:
                pass
            return None, None

        def apply_cell_format(cell, font_ref, align_ref):
            if not cell.text_frame.paragraphs:
                return
            try:
                cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
            tgt = cell.text_frame.paragraphs[0]
            if font_ref:
                _apply_paragraph_font(tgt, font_ref)
            if align_ref is not None:
                try:
                    tgt.alignment = align_ref
                except Exception:
                    pass

        # 헤더 행 (0행)
        for c in range(len(headers)):
            col = data_start_col + c
            if col >= ncols:
                break
            cell = tbl.cell(0, col)
            font_ref, align_ref = cell_format_ref(0, col)
            cell.text = headers[c] if c < len(headers) else ""
            apply_cell_format(cell, font_ref, align_ref)

        # 템플릿 행 수가 부족하면 마지막 행을 XML 복제해 행 추가
        needed = 1 + len(rows)  # 헤더 + 데이터
        if needed > nrows_total:
            try:
                # NOTE: _tbl은 python-pptx 비공개 속성(lxml Element)으로,
                # 공식 API에 행 추가 메서드가 없어 직접 XML을 조작한다.
                # python-pptx 버전 업데이트 시 동작 확인 필요.
                tbl_el = tbl._tbl
                ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                tr_list = tbl_el.findall(f'{ns}tr')
                last_tr = tr_list[-1] if tr_list else None
                if last_tr is not None:
                    for _ in range(needed - nrows_total):
                        new_tr = copy.deepcopy(last_tr)
                        tbl_el.append(new_tr)
                    nrows_total = needed
            except Exception as e:
                print(f"경고: 테이블 행 추가 실패 (프로퍼티 {len(rows)}개 중 {nrows_total - 1}개만 표시): {e}", file=sys.stderr)

        # 데이터 행
        for r, row in enumerate(rows):
            if r + 1 >= nrows_total:
                break
            if data_start_col == 1:
                try:
                    cell_no = tbl.cell(r + 1, 0)
                    font_ref, align_ref = cell_format_ref(r + 1, 0)
                    cell_no.text = str(r + 1)
                    apply_cell_format(cell_no, font_ref, align_ref)
                except Exception:
                    pass
            for c in range(len(headers)):
                col = data_start_col + c
                if col >= ncols:
                    break
                cell = tbl.cell(r + 1, col)
                font_ref, align_ref = cell_format_ref(r + 1, col)
                cell.text = str(row[c]) if c < len(row) else ""
                apply_cell_format(cell, font_ref, align_ref)

        # 우리 데이터 행 밑 비우기 + 해당 셀 서식 유지
        for r in range(1 + len(rows), nrows_total):
            for c in range(ncols):
                try:
                    cell = tbl.cell(r, c)
                    font_ref, align_ref = cell_format_ref(r, c)
                    cell.text = ""
                    apply_cell_format(cell, font_ref, align_ref)
                except Exception:
                    pass
    except Exception as e:
        print(f"경고: 테이블 채우기 중 오류 {e}", file=sys.stderr)


def main():
    if len(sys.argv) < 3:
        print("사용법: python3 generate_pptx.py <데이터.md> <출력.pptx> [템플릿.pptx]", file=sys.stderr)
        sys.exit(1)
    data_md = Path(sys.argv[1])
    output_pptx = Path(sys.argv[2])
    template_pptx = Path(sys.argv[3]) if len(sys.argv) > 3 else DEFAULT_TEMPLATE

    sys.path.insert(0, str(SCRIPT_DIR))
    from parse_data_md import parse_md

    data = parse_md(data_md)
    meta = data["meta"]
    screens = data["태깅 대상 화면"]
    events = data["이벤트"]

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("python-pptx가 필요합니다: pip install python-pptx", file=sys.stderr)
        sys.exit(2)

    if template_pptx.exists():
        prs = Presentation(str(template_pptx))
        detail_template_idx = 3
        if len(events) == 0:
            # 이벤트가 없으면 태깅 상세 슬라이드(인덱스 3) 제거
            # NOTE: _sldIdLst는 python-pptx 비공개 속성으로, 공식 API에
            # 슬라이드 삭제 메서드가 없어 직접 조작한다.
            # python-pptx 버전 업데이트 시 동작 확인 필요.
            if len(prs.slides) > detail_template_idx:
                rId = prs.slides._sldIdLst[detail_template_idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[detail_template_idx])
        else:
            # 태깅 상세 슬라이드를 이벤트 수만큼 복제
            for _ in range(len(events) - 1):
                _duplicate_slide(prs, detail_template_idx)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

    def set_shape_text(slide, text: str, idx: int = 0):
        """지정 인덱스의 텍스트 shape에 제목 설정. runs[0].text만 바꾸고 나머지 run은 비워 폰트/서식 유지."""
        for i, s in enumerate(slide.shapes):
            if s.has_text_frame and i == idx:
                tf = s.text_frame
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    if p.runs:
                        p.runs[0].text = text
                        for j in range(1, len(p.runs)):
                            p.runs[j].text = ""
                    else:
                        p.text = text
                return
        left, top, width, height = Inches(0.5), Inches(1), Inches(9), Inches(1)
        tx = slide.shapes.add_textbox(left, top, width, height)
        tx.text_frame.paragraphs[0].text = text
        tx.text_frame.paragraphs[0].font.size = Pt(24)

    def add_table(slide, headers: list, rows: list, left_inches: float = 0.5, width_inches: float = 9.0):
        from pptx.util import Inches
        rows_total = 1 + len(rows)
        cols = len(headers)
        left, width = Inches(left_inches), Inches(width_inches)
        top, height = Inches(1.5), Inches(0.25 * (rows_total + 1))
        table_shape = slide.shapes.add_table(rows_total, cols, left, top, width, height)
        table = table_shape.table
        for c, h in enumerate(headers):
            table.cell(0, c).text = h
        for r, row in enumerate(rows):
            for c in range(len(headers)):
                table.cell(r + 1, c).text = str(row[c]) if c < len(row) else ""
        return table_shape

    def add_picture(slide, image_path: Path, left_inches: float = 0.5, top_inches: float = 1.5,
                    width_inches: float = WIREFRAME_WIDTH_INCH, max_height_inches: float = WIREFRAME_HEIGHT_INCH):
        from pptx.util import Inches
        try:
            from PIL import Image as PILImage
            with PILImage.open(str(image_path)) as img:
                img_w, img_h = img.size
            aspect = img_h / img_w  # 세로/가로 비율

            # 가로 기준으로 height 계산 후 영역 초과 여부 판단
            calc_height = width_inches * aspect
            if calc_height > max_height_inches:
                # 모바일 세로 캡처처럼 세로가 긴 경우: height 기준으로 맞춤
                final_height = max_height_inches
                final_width = max_height_inches / aspect
            else:
                # PC 가로 캡처처럼 가로가 넓은 경우: width 기준으로 맞춤
                final_width = width_inches
                final_height = calc_height

            slide.shapes.add_picture(
                str(image_path),
                Inches(left_inches), Inches(top_inches),
                Inches(final_width), Inches(final_height),
            )
        except Exception as e:
            print(f"경고: 이미지 삽입 실패 {image_path}: {e}", file=sys.stderr)

    # --- 표지: 제목 + 작성일/팀/작성자. 테이블이 있으면 값 열만 채움; 없으면 placeholder별로 한 칸에 한 값(shape[1]=작성일, [2]=팀, [3]=작성자) ---
    if len(prs.slides) >= 1:
        s0 = prs.slides[0]
        set_shape_text(s0, f"{meta['서비스명']} 태깅 정의서", 0)
        cover_table_found = False
        for sh in s0.shapes:
            if sh.has_table:
                tbl = sh.table
                if len(tbl.rows) >= 1 and len(tbl.rows[0].cells) >= 2:
                    for r in range(len(tbl.rows)):
                        try:
                            label_cell = (tbl.cell(r, 0).text or "").strip()
                            if "소속" in label_cell:
                                tbl.cell(r, 1).text = meta["팀"]
                            elif "작성자" in label_cell and "최종" not in label_cell:
                                tbl.cell(r, 1).text = meta["작성자"]
                            elif "최종" in label_cell or "업데이트" in label_cell or "작성일" in label_cell:
                                tbl.cell(r, 1).text = meta["작성일"]
                        except IndexError:
                            pass
                    cover_table_found = True
                break
        if not cover_table_found:
            # placeholder별 한 칸에 한 값만 (폰트/서식 유지 위해 첫 단락만 설정)
            vals = [meta["작성일"], meta["팀"], meta["작성자"]]
            for vi, idx in enumerate([1, 2, 3]):
                if idx < len(s0.shapes) and s0.shapes[idx].has_text_frame and s0.shapes[idx].text_frame.paragraphs:
                    s0.shapes[idx].text_frame.paragraphs[0].text = vals[vi] if vi < len(vals) else ""

    # --- 이력: 버전 0.0.1, 변경일=표지 일자, 변경내용=표지제목+초안, 작성자=표지 작성자. 헤더 행이 있으면 데이터만 1행에 채움, 2열이면 4행 채움 ---
    if len(prs.slides) >= 2:
        s1 = prs.slides[1]
        set_shape_text(s1, "HISTORY", 0)
        # HISTORY 제목 폰트 크기 60% 축소
        if s1.shapes[0].has_text_frame and s1.shapes[0].text_frame.paragraphs:
            p0 = s1.shapes[0].text_frame.paragraphs[0]
            finfo = _get_paragraph_font(p0)
            if finfo.get("size"):
                new_size = int(finfo["size"] * 0.4)
                _apply_paragraph_font(p0, {"size": new_size})
        version_val = "0.0.1"
        date_val = meta["작성일"]
        change_val = f"{meta['서비스명']} 태깅 정의 초안"
        author_val = meta["작성자"]
        history_lines = [
            f"버전: {version_val}",
            f"변경일: {date_val}",
            f"변경내용: {change_val}",
            f"작성자: {author_val}",
        ]
        table_found = False
        for sh in s1.shapes:
            if sh.has_table:
                tbl = sh.table
                nrows, ncols = len(tbl.rows), len(tbl.rows[0].cells) if tbl.rows else 0
                if nrows >= 1 and ncols >= 2:
                    r0c0 = (tbl.cell(0, 0).text or "").strip()
                    r0c1 = (tbl.cell(0, 1).text or "").strip()
                    # 헤더 행 여부: 첫 행에 "버전","변경일" 등이 컬럼으로 있으면 데이터는 1행에 채움
                    is_header_row = ncols >= 4 and ("변경일" in (tbl.cell(0, 1).text or "") or "변경내용" in (tbl.cell(0, 2).text or ""))
                    if is_header_row and nrows >= 2:
                        # 1행 셀별로 템플릿 폰트·정렬 유지: 덮어쓰기 전에 서식 읽고, 텍스트 설정 후 적용
                        row1_vals = [version_val, date_val, change_val, author_val, ""]
                        for c in range(ncols):
                            try:
                                cell = tbl.cell(1, c)
                                font_ref = _get_paragraph_font(cell.text_frame.paragraphs[0]) if cell.text_frame.paragraphs else None
                                cell.text = row1_vals[c] if c < len(row1_vals) else ""
                                if font_ref and cell.text_frame.paragraphs:
                                    _apply_paragraph_font(cell.text_frame.paragraphs[0], font_ref)
                            except Exception:
                                pass
                        # 데이터 행 외(2행~) 비우기 — 셀 서식은 유지
                        for r in range(2, nrows):
                            for c in range(ncols):
                                try:
                                    cell = tbl.cell(r, c)
                                    font_ref = _get_paragraph_font(cell.text_frame.paragraphs[0]) if cell.text_frame.paragraphs else None
                                    cell.text = ""
                                    if font_ref and cell.text_frame.paragraphs:
                                        _apply_paragraph_font(cell.text_frame.paragraphs[0], font_ref)
                                except Exception:
                                    pass
                        # 마지막 행 1개 제거 (표 사이즈 축소)
                        try:
                            tbl_el = tbl._tbl
                            tr_list = tbl_el.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}tr')
                            if len(tr_list) > 2:
                                tbl_el.remove(tr_list[-1])
                        except Exception:
                            pass
                        table_found = True
                        break
                    # 2열(라벨|값) 형태: 4행 채움 — 셀별 폰트 유지
                    if nrows >= 4 and ncols >= 2:
                        pairs = [
                            ("버전", version_val), ("변경일", date_val),
                            ("변경내용", change_val), ("작성자", author_val),
                        ]
                        for r in range(min(4, nrows)):
                            for c in range(2):
                                try:
                                    cell = tbl.cell(r, c)
                                    font_ref = _get_paragraph_font(cell.text_frame.paragraphs[0]) if cell.text_frame.paragraphs else None
                                    cell.text = pairs[r][c] if r < len(pairs) else ""
                                    if font_ref and cell.text_frame.paragraphs:
                                        _apply_paragraph_font(cell.text_frame.paragraphs[0], font_ref)
                                except Exception:
                                    pass
                        table_found = True
                        break
                break
        if not table_found and len(s1.shapes) > 1 and s1.shapes[1].has_text_frame:
            _set_text_frame_lines(s1.shapes[1], history_lines)

    # --- 개요: 타이틀="개요"(shape[1]), 본문="태깅 범위"+블릿(shape[2]만 채우고 [3] 등 비움) — 폰트 유지 ---
    if len(prs.slides) >= 3:
        s2 = prs.slides[2]
        if len(s2.shapes) > 1 and s2.shapes[1].has_text_frame:
            s2.shapes[1].text_frame.paragraphs[0].text = "개요"
            # 제목 폰트 유지(이미 같은 shape라 텍스트만 바꿨으면 서식 유지)
        body_lines = ["태깅 범위"]
        if screens:
            body_lines.extend(f"- {x}" for x in screens)
        if len(s2.shapes) > 2 and s2.shapes[2].has_text_frame:
            _set_text_frame_lines(s2.shapes[2], body_lines, preserve_font_from_shape=s2.shapes[2])
            # 본문 줄간격 150%
            try:
                from pptx.util import Pt
                from pptx.oxml.ns import qn
                for para in s2.shapes[2].text_frame.paragraphs:
                    pPr = para._p.get_or_add_pPr()
                    lnSpc = pPr.find(qn("a:lnSpc"))
                    if lnSpc is None:
                        from pptx.oxml.xmlchemy import OxmlElement
                        lnSpc = OxmlElement("a:lnSpc")
                        pPr.append(lnSpc)
                    else:
                        for child in list(lnSpc):
                            lnSpc.remove(child)
                    from pptx.oxml.xmlchemy import OxmlElement
                    spcPct = OxmlElement("a:spcPct")
                    spcPct.set("val", "150000")  # 150%
                    lnSpc.append(spcPct)
            except Exception:
                pass
        for idx in range(3, len(s2.shapes)):
            if s2.shapes[idx].has_text_frame:
                for p in s2.shapes[idx].text_frame.paragraphs:
                    p.text = ""

    # --- 태깅 상세: 모든 슬라이드가 템플릿 복제본이므로 데이터만 교체 ---
    detail_start = 3
    base_dir = data_md.resolve().parent
    project_root = _find_project_root(data_md)
    headers = ["Event property", "property 구분", "value", "description"]

    for i, ev in enumerate(events):
        slide = prs.slides[detail_start + i]

        # 1) 제목 교체: shape[1] = 제목 TEXT_BOX (runs[0].text만 바꿔 폰트 유지)
        title = ev.get("Event Name (한글)") or ev.get("Event Name (영문)") or f"이벤트 {i+1}"
        title_str = f"{title}  ({ev.get('호출 시점', '')})"
        set_shape_text(slide, title_str, 1)

        # 2) 이미지 교체: 기존 PICTURE 제거 후 캡처 삽입 (없으면 와이어프레임)
        rows = [[p["Event property"], p["property 구분"], p["value"], p["description"]] for p in ev["properties"]]
        img_path_raw = (ev.get("캡처 경로") or "").strip()
        img_path = _resolve_image_path(img_path_raw, base_dir, project_root)
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for s in list(slide.shapes):
                if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    s._element.getparent().remove(s._element)
        except Exception:
            pass
        if img_path and img_path.exists():
            add_picture(slide, img_path, left_inches=WIREFRAME_LEFT_INCH, top_inches=WIREFRAME_TOP_INCH, width_inches=WIREFRAME_WIDTH_INCH)
        else:
            _add_wireframe_placeholder(slide)

        # 3) 표 데이터 채우기: 복제된 표에 데이터만 교체 (서식 그대로 유지)
        if not rows:
            continue
        table_shape = None
        for sh in slide.shapes:
            if sh.has_table:
                table_shape = sh
                break
        if table_shape:
            _fill_table_from_template(table_shape.table, headers, rows, copy_style_from=table_shape.table)
            _apply_tagging_table_overrides_only(table_shape)

    # --- 전체 폰트 통일: 나눔고딕, 볼드는 표지 제목 / 각 장표 타이틀 / 표 헤더에만 ---
    FONT_NAME = "나눔고딕"
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # 표 처리
            if shape.has_table:
                tbl = shape.table
                for r in range(len(tbl.rows)):
                    for c in range(len(tbl.rows[0].cells)):
                        try:
                            for para in tbl.cell(r, c).text_frame.paragraphs:
                                is_header = (r == 0)
                                for run in para.runs:
                                    run.font.name = FONT_NAME
                                    run.font.bold = is_header
                                if not para.runs:
                                    para.font.name = FONT_NAME
                                    para.font.bold = is_header
                        except Exception:
                            pass
                continue
            if not shape.has_text_frame:
                continue
            # 타이틀 판별: 표지(si==0) shape[0], 이력(si==1) shape[0],
            # 개요(si==2) shape[1], 태깅 상세(si>=3) shape[1]
            is_title_shape = False
            shape_idx = list(slide.shapes).index(shape)
            if si == 0 and shape_idx == 0:
                is_title_shape = True
            elif si == 1 and shape_idx == 0:
                is_title_shape = True
            elif si == 2 and shape_idx == 1:
                is_title_shape = True
            elif si >= 3 and shape_idx == 1:
                is_title_shape = True
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT_NAME
                    run.font.bold = is_title_shape
                if not para.runs:
                    para.font.name = FONT_NAME
                    para.font.bold = is_title_shape

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    print(f"저장: {output_pptx}")


if __name__ == "__main__":
    main()
