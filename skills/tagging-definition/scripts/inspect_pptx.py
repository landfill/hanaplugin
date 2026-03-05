#!/usr/bin/env python3
"""템플릿/생성된 PPT의 슬라이드별 shape·폰트·테이블 구조를 텍스트로 덤프."""
import sys
from pathlib import Path

def _font_info(shape):
    try:
        if shape.has_text_frame and shape.text_frame.paragraphs:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                r = p.runs[0]
                return f"size={r.font.size}, bold={r.font.bold}, name={getattr(r.font.name, '__str__', lambda: r.font.name)()}"
            if hasattr(p, 'font') and p.font:
                return f"size={p.font.size}, bold={p.font.bold}"
    except Exception as e:
        return str(e)
    return ""

def _cell_font(cell):
    try:
        if cell.text_frame.paragraphs:
            p = cell.text_frame.paragraphs[0]
            if p.runs:
                r = p.runs[0]
                return f"size={r.font.size}"
    except Exception:
        pass
    return ""

def inspect_pptx(pptx_path: Path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(str(pptx_path))
    lines = []
    for si, slide in enumerate(prs.slides):
        lines.append(f"\n=== 슬라이드 {si + 1} (인덱스 {si}) ===")
        for idx, shape in enumerate(slide.shapes):
            stype = getattr(shape, "shape_type", None)
            type_name = str(stype) if stype is not None else type(shape).__name__
            lines.append(f"  [{idx}] type={type_name} left={getattr(shape, 'left', None)} top={getattr(shape, 'top', None)}")
            if shape.has_text_frame:
                full = []
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    t = (para.text or "").strip()
                    if t:
                        full.append(t[:60])
                if full:
                    lines.append(f"       text_frame: {repr(full)}")
                    lines.append(f"       font(첫단락): {_font_info(shape)}")
            if shape.has_table:
                tbl = shape.table
                rows, cols = len(tbl.rows), len(tbl.rows[0].cells) if tbl.rows else 0
                lines.append(f"       table: {rows}행 x {cols}열")
                for r in range(min(rows, 6)):
                    row_texts = []
                    for c in range(min(cols, 5)):
                        cell = tbl.cell(r, c)
                        row_texts.append((cell.text or "").strip()[:20])
                    lines.append(f"         row{r}: {row_texts}")
                    if r == 0 and cols > 0:
                        lines.append(f"         row0 font: {_cell_font(tbl.cell(0, 0))}")
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                lines.append(f"       [PICTURE]")
    return "\n".join(lines)

if __name__ == "__main__":
    p = Path(sys.argv[1])
    print(inspect_pptx(p))
